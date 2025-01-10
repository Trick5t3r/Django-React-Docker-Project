from django.db import models
from django.contrib.auth.models import User

from django.utils.timezone import now
import os
import textract
import PyPDF2
import openpyxl
import pptx
from moviepy import VideoFileClip
import speech_recognition as sr
from PIL import Image
#import pytesseract
from .ai_files.llm_tools import generate_title
from mistralai import Mistral
from dotenv import load_dotenv
from polymorphic.models import PolymorphicModel


# Create your models here.

class ChatSession(models.Model):
    title = models.CharField(max_length=100)
    user = models.ForeignKey(User, on_delete=models.CASCADE, related_name="sessions")
    messages = models.JSONField(default=list) 
    created_at = models.DateTimeField(auto_now_add=True)

    def __str__(self):
        return f"Session with {self.user.username} at {self.created_at}"
    


###### RAG Models

load_dotenv()

# Configuration du client Mistral
model = "mistral-embedding-latest"  # Modèle d'embedding adapté
api_key = os.environ.get("MISTRAL_API_KEY")
client = Mistral(api_key=api_key)

# Base Model
class BaseFile(PolymorphicModel):
    title = models.CharField(blank=True, null=True, max_length=255)
    chatsession = models.ForeignKey(ChatSession, on_delete=models.CASCADE)
    created_at = models.DateTimeField(default=now)
    file = models.FileField()#upload_to='uploads/')
    content = models.TextField(blank=True, null=True)  # For storing extracted text
    embeddings = models.JSONField(blank=True, null=True)  # For storing vector embeddings

    def extract_data(self):
        """Default extract method for text files."""
        try:
            self.content = textract.process(self.file.path).decode('utf-8')
        except Exception as e:
            self.content = f"Error extracting data: {str(e)}"
        self.index_content()
        return self.embeddings

    def index_content(self):
        """Generate embeddings and index content."""
        if self.content:
            self.embeddings = self.generate_embeddings(self.content)

    def generate_embeddings(self, content):
        """
        Generate embeddings using Mistral's API.

        Args:
            content (str): The text content to generate embeddings for.

        Returns:
            list: A list of embedding values (float).
        """
        try:
            # Envoyer le contenu à l'API Mistral pour obtenir les embeddings
            response = client.embeddings.create(
                model=model,
                input=content[:1000],  # Limite à 1000 caractères pour s'assurer de respecter les contraintes
            )

            # Récupérer les embeddings
            print(response)
            embeddings = response.embeddings
            return embeddings
        except Exception as e:
            return []#ValueError(f"Error generating embeddings: {str(e)}")


# PDF File Model
class PDFFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            with open(self.file.path, 'rb') as pdf_file:
                reader = PyPDF2.PdfReader(pdf_file)
                self.content = " ".join(page.extract_text() for page in reader.pages)
        except Exception as e:
            self.content = f"Error extracting data from PDF: {str(e)}"
        self.index_content()
        return self.embeddings


# Excel File Model
class ExcelFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            workbook = openpyxl.load_workbook(self.file.path)
            sheet = workbook.active
            self.content = "\n".join(
                ["\t".join([str(cell.value) for cell in row]) for row in sheet.iter_rows()]
            )
        except Exception as e:
            self.content = f"Error extracting data from Excel: {str(e)}"
        self.index_content()
        return self.embeddings


# PowerPoint File Model
class PowerPointFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            presentation = pptx.Presentation(self.file.path)
            self.content = "\n".join(
                [paragraph.text for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs]
            )
        except Exception as e:
            self.content = f"Error extracting data from PowerPoint: {str(e)}"
        self.index_content()
        return self.embeddings


# Audio File Model
class AudioFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            recognizer = sr.Recognizer()
            with sr.AudioFile(self.file.path) as source:
                audio_data = recognizer.record(source)
                self.content = recognizer.recognize_google(audio_data)
        except Exception as e:
            self.content = f"Error extracting data from audio file: {str(e)}"
        self.index_content()
        return self.embeddings


# Video File Model
class VideoFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            video = VideoFileClip(self.file.path)
            self.content = f"Video file with duration {video.duration} seconds."
        except Exception as e:
            self.content = f"Error extracting data from video file: {str(e)}"
        self.index_content()
        return self.embeddings


# Image File Model
class ImageFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first

    def extract_data(self):
        try:
            # Utiliser Pillow et Tesseract pour extraire du texte d'une image
            image = Image.open(self.file)
            self.content = "" #pytesseract.image_to_string(image, lang='eng')
        except Exception as e:
            self.content = f"Error extracting data from image file: {str(e)}"
        self.index_content()
        return self.embeddings


# Other File Model
class OtherFile(BaseFile):
    def save(self, *args, **kwargs):
        """Override save to determine and create specific subclass."""
        if not self.content or not self.embeddings:
            self.extract_data()
        if self.file and not self.title:  # Check if the file exists and title is not set
            self.title = os.path.basename(self.file.name)  # Extract the file name
        super().save(*args, **kwargs)  # Save the file first
